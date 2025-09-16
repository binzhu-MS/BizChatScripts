#!/usr/bin/env python3
"""
Universal File Reader Utility
============================

This module provides comprehensive file reading capabilities for multiple document formats
including Word, PDF, Excel, and PowerPoint files. It uses multiple fallback methods for
each format to ensure maximum compatibility and reliability.

Supported Formats:
- PDF files (.pdf) - Using PyPDF2 with fallback methods
- Word documents (.docx, .doc) - Using python-docx and docx2txt
- Excel files (.xlsx, .xls, .xlsb) - Using pandas, openpyxl, xlrd, pyxlsb, and COM automation
- PowerPoint files (.pptx, .ppt) - Using python-pptx and COM automation
- Text files (.txt) - Direct reading with encoding detection

Key Features:
- Multiple fallback methods for maximum file compatibility
- Comprehensive error handling and logging
- Structured content extraction preserving document organization
- Support for legacy and modern file formats
- Windows COM automation as fallback for Office files
- Fire CLI integration for command-line usage

Usage:
    python universal_file_reader.py read_file --file_path="document.pdf"
    python universal_file_reader.py get_file_info --file_path="spreadsheet.xlsx"
    python universal_file_reader.py test_capabilities

    # In Python code:
    from universal_file_reader import UniversalFileReader
    reader = UniversalFileReader()
    content = reader.read_file("document.pdf")
"""

import os
import sys
import logging
import fire
import chardet
from pathlib import Path
from typing import Dict, List, Optional, Any, Union, Tuple
from datetime import datetime

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] [%(name)s] %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
logger = logging.getLogger(__name__)


class UniversalFileReader:
    """Universal file reader for multiple document formats with comprehensive fallback methods"""

    def __init__(self, silent: bool = False):
        """Initialize the universal file reader

        Args:
            silent: If True, suppress info logging (only warnings and errors will be shown)
        """
        self.silent = silent
        self.supported_formats = {
            "pdf": [".pdf"],
            "word": [".docx", ".doc"],
            "excel": [".xlsx", ".xls", ".xlsb"],
            "powerpoint": [".pptx", ".ppt"],
            "text": [".txt", ".md", ".rtf"],
        }
        self.all_extensions = []
        for extensions in self.supported_formats.values():
            self.all_extensions.extend(extensions)

        logger.debug("📚 UniversalFileReader initialized")

    def _log_info(self, message: str) -> None:
        """Log info message only if not in silent mode"""
        if not self.silent:
            logger.info(message)

    def read_file(self, file_path: str) -> Dict[str, Any]:
        """Read content from any supported file format with comprehensive error handling

        Args:
            file_path: Path to the file to read

        Returns:
            Dict[str, Any]: File content and metadata

        Raises:
            FileNotFoundError: If the file doesn't exist
            ValueError: If the file format is not supported
        """
        self._log_info(f"📖 Reading file: {file_path}")

        # Validate file exists
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"❌ File not found: {file_path}")

        # Get file info
        file_info = self._get_file_info(file_path)
        file_extension = file_info["file_extension"]

        # Determine file type
        file_type = self._detect_file_type(file_extension)

        if not file_type:
            raise ValueError(f"❌ Unsupported file format: {file_extension}")

        self._log_info(f"🔍 Detected file type: {file_type}")

        try:
            # Route to appropriate reader
            if file_type == "pdf":
                content_data = self._read_pdf(file_path)
            elif file_type == "word":
                content_data = self._read_word(file_path)
            elif file_type == "excel":
                content_data = self._read_excel(file_path)
            elif file_type == "powerpoint":
                content_data = self._read_powerpoint(file_path)
            elif file_type == "text":
                content_data = self._read_text(file_path)
            else:
                raise ValueError(f"❌ Unsupported file type: {file_type}")

            # Combine file info with content
            result = {
                **file_info,
                **content_data,
                "read_timestamp": datetime.now().isoformat(),
                "status": "success",
                "format": file_type,  # Add format type
                "method": content_data.get(
                    "extraction_method", "unknown"
                ),  # Standardize method name
            }

            self._log_info(
                f"✅ Successfully read file: {len(result.get('content', ''))} characters extracted"
            )
            return result

        except Exception as e:
            error_result = {
                **file_info,
                "content": "",
                "error": str(e),
                "read_timestamp": datetime.now().isoformat(),
                "status": "error",
            }
            logger.error(f"❌ Failed to read file {file_path}: {e}")
            return error_result

    def _detect_file_type(self, extension: str) -> Optional[str]:
        """Detect file type from extension

        Args:
            extension: File extension (with dot)

        Returns:
            Optional[str]: File type or None if not supported
        """
        extension = extension.lower()
        for file_type, extensions in self.supported_formats.items():
            if extension in extensions:
                return file_type
        return None

    def _get_file_info(self, file_path: str) -> Dict[str, Any]:
        """Get comprehensive file information

        Args:
            file_path: Path to the file

        Returns:
            Dict[str, Any]: File metadata
        """
        file_path_obj = Path(file_path)
        file_stat = file_path_obj.stat()

        return {
            "file_name": file_path_obj.name,
            "file_path": str(file_path_obj.absolute()),
            "file_extension": file_path_obj.suffix.lower(),
            "file_size_bytes": file_stat.st_size,
            "file_size_mb": round(file_stat.st_size / (1024 * 1024), 2),
            "modified_time": datetime.fromtimestamp(file_stat.st_mtime).isoformat(),
        }

    def _read_pdf(self, file_path: str) -> Dict[str, Any]:
        """Read PDF files with multiple fallback methods

        Args:
            file_path: Path to PDF file

        Returns:
            Dict[str, Any]: PDF content and metadata
        """
        self._log_info(f"📄 Reading PDF: {file_path}")

        # Method 1: PyPDF2
        try:
            import PyPDF2

            content = ""
            pages_processed = 0

            with open(file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                total_pages = len(pdf_reader.pages)

                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            content += f"\n--- Page {page_num + 1} ---\n"
                            content += page_text + "\n"
                        pages_processed += 1
                    except Exception as e:
                        logger.warning(f"⚠️ Error reading page {page_num + 1}: {e}")
                        continue

            self._log_info(
                f"✅ PyPDF2: {pages_processed}/{total_pages} pages processed"
            )
            return {
                "content": content.strip(),
                "total_pages": total_pages,
                "pages_processed": pages_processed,
                "extraction_method": "PyPDF2",
            }

        except ImportError:
            logger.warning("⚠️ PyPDF2 not available")
        except Exception as e:
            logger.warning(f"⚠️ PyPDF2 failed: {e}")

        # Method 2: Try other PDF libraries (if available)
        try:
            import fitz  # PyMuPDF

            content = ""
            doc = fitz.open(file_path)
            total_pages = doc.page_count

            for page_num in range(total_pages):
                try:
                    page = doc[page_num]
                    page_text = page.get_text()
                    if page_text.strip():
                        content += f"\n--- Page {page_num + 1} ---\n"
                        content += page_text + "\n"
                except Exception as e:
                    logger.warning(f"⚠️ Error reading page {page_num + 1}: {e}")
                    continue

            doc.close()
            self._log_info(f"✅ PyMuPDF: {total_pages} pages processed")
            return {
                "content": content.strip(),
                "total_pages": total_pages,
                "pages_processed": total_pages,
                "extraction_method": "PyMuPDF",
            }

        except ImportError:
            logger.warning("⚠️ PyMuPDF not available")
        except Exception as e:
            logger.warning(f"⚠️ PyMuPDF failed: {e}")

        raise Exception("All PDF reading methods failed")

    def _read_word(self, file_path: str) -> Dict[str, Any]:
        """Read Word documents with multiple fallback methods

        Args:
            file_path: Path to Word document

        Returns:
            Dict[str, Any]: Word document content and metadata
        """
        self._log_info(f"📝 Reading Word document: {file_path}")

        file_extension = Path(file_path).suffix.lower()

        if file_extension == ".docx":
            return self._read_docx(file_path)
        elif file_extension == ".doc":
            return self._read_doc(file_path)
        else:
            raise ValueError(f"Unsupported Word format: {file_extension}")

    def _read_docx(self, file_path: str) -> Dict[str, Any]:
        """Read .docx files

        Args:
            file_path: Path to .docx file

        Returns:
            Dict[str, Any]: Document content and metadata
        """
        # Method 1: python-docx
        try:
            from docx import Document

            doc = Document(file_path)
            content = ""
            paragraphs_count = 0
            tables_count = 0

            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content += paragraph.text + "\n"
                    paragraphs_count += 1

            # Extract tables
            for table in doc.tables:
                content += "\n--- Table ---\n"
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_data.append(cell_text)
                    if row_data:
                        content += " | ".join(row_data) + "\n"
                tables_count += 1

            self._log_info(
                f"✅ python-docx: {paragraphs_count} paragraphs, {tables_count} tables"
            )
            return {
                "content": content.strip(),
                "paragraphs_count": paragraphs_count,
                "tables_count": tables_count,
                "extraction_method": "python-docx",
            }

        except ImportError:
            logger.warning("⚠️ python-docx not available")
        except Exception as e:
            logger.warning(f"⚠️ python-docx failed: {e}")

        # Method 2: docx2txt
        try:
            import docx2txt

            content = docx2txt.process(file_path)
            content = content.strip() if content else ""

            self._log_info(f"✅ docx2txt: {len(content)} characters extracted")
            return {"content": content, "extraction_method": "docx2txt"}

        except ImportError:
            logger.warning("⚠️ docx2txt not available")
        except Exception as e:
            logger.warning(f"⚠️ docx2txt failed: {e}")

        raise Exception("All .docx reading methods failed")

    def _read_doc(self, file_path: str) -> Dict[str, Any]:
        """Read .doc files

        Args:
            file_path: Path to .doc file

        Returns:
            Dict[str, Any]: Document content and metadata
        """
        # Method 1: docx2txt (handles .doc files)
        try:
            import docx2txt

            content = docx2txt.process(file_path)
            content = content.strip() if content else ""

            self._log_info(f"✅ docx2txt: {len(content)} characters extracted")
            return {"content": content, "extraction_method": "docx2txt"}

        except ImportError:
            logger.warning("⚠️ docx2txt not available")
        except Exception as e:
            logger.warning(f"⚠️ docx2txt failed: {e}")

        # Method 2: Windows COM (if on Windows)
        try:
            import win32com.client
            import pythoncom

            pythoncom.CoInitialize()
            try:
                word_app = win32com.client.DispatchEx("Word.Application")
                word_app.Visible = False

                doc = word_app.Documents.Open(os.path.abspath(file_path), ReadOnly=True)
                content = doc.Content.Text
                doc.Close(SaveChanges=False)
                word_app.Application.Quit()

                self._log_info(f"✅ COM Word: {len(content)} characters extracted")
                return {"content": content.strip(), "extraction_method": "COM Word"}

            finally:
                pythoncom.CoUninitialize()

        except ImportError:
            logger.warning("⚠️ pywin32 not available")
        except Exception as e:
            logger.warning(f"⚠️ COM Word failed: {e}")

        raise Exception("All .doc reading methods failed")

    def _read_excel(self, file_path: str) -> Dict[str, Any]:
        """Read Excel files with comprehensive fallback methods (based on excel_to_tsv.py)

        Args:
            file_path: Path to Excel file

        Returns:
            Dict[str, Any]: Excel content and metadata
        """
        self._log_info(f"📊 Reading Excel file: {file_path}")

        # Get sheet names first
        sheet_names = self._get_excel_sheets(file_path)

        if not sheet_names:
            raise Exception("Could not read Excel file or file contains no sheets")

        # Read all sheets
        all_content = ""
        sheets_data = {}

        for sheet_name in sheet_names:
            try:
                sheet_content = self._read_excel_sheet(file_path, sheet_name)
                sheets_data[sheet_name] = sheet_content
                all_content += f"\n--- Sheet: {sheet_name} ---\n"
                all_content += sheet_content + "\n"
            except Exception as e:
                logger.warning(f"⚠️ Error reading sheet '{sheet_name}': {e}")
                continue

        self._log_info(f"✅ Excel: {len(sheet_names)} sheets processed")
        return {
            "content": all_content.strip(),
            "sheet_names": sheet_names,
            "sheets_data": sheets_data,
            "total_sheets": len(sheet_names),
            "extraction_method": "pandas + fallbacks",
        }

    def _get_excel_sheets(self, file_path: str) -> List[str]:
        """Get sheet names using multiple fallback methods (adapted from excel_to_tsv.py)

        Args:
            file_path: Path to Excel file

        Returns:
            List[str]: Sheet names
        """
        # Check file header to determine format
        with open(file_path, "rb") as f:
            header = f.read(4)

        engines_to_try = []

        if header.startswith(b"PK"):
            # ZIP-based format (xlsx, xlsm)
            engines_to_try = ["openpyxl", "xlrd", None]
        elif header.startswith(b"\xd0\xcf\x11\xe0"):
            # OLE2 format (xls)
            engines_to_try = ["xlrd", "openpyxl", None]
        else:
            # Unknown format, try all
            engines_to_try = ["openpyxl", "xlrd", None]

        # Try pandas with different engines
        for engine in engines_to_try:
            try:
                import pandas as pd

                engine_name = engine if engine else "default"
                logger.debug(f"Trying pandas with {engine_name} engine...")

                with pd.ExcelFile(file_path, engine=engine) as xls:
                    sheets = [str(sheet) for sheet in xls.sheet_names]
                    self._log_info(
                        f"✅ {engine_name} engine: {len(sheets)} sheets found"
                    )
                    return sheets

            except Exception as e:
                logger.debug(f"❌ {engine_name} engine failed: {e}")
                continue

        # Try direct library access
        try:
            import openpyxl

            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            sheets = [str(sheet) for sheet in wb.sheetnames]
            wb.close()
            self._log_info(f"✅ Direct openpyxl: {len(sheets)} sheets found")
            return sheets
        except Exception as e:
            logger.debug(f"❌ Direct openpyxl failed: {e}")

        try:
            import xlrd

            wb = xlrd.open_workbook(file_path, ignore_workbook_corruption=True)
            sheets = [str(sheet) for sheet in wb.sheet_names()]
            self._log_info(f"✅ Direct xlrd: {len(sheets)} sheets found")
            return sheets
        except Exception as e:
            logger.debug(f"❌ Direct xlrd failed: {e}")

        # Windows COM as last resort
        try:
            import win32com.client
            import pythoncom

            pythoncom.CoInitialize()
            try:
                excel_app = win32com.client.DispatchEx("Excel.Application")
                excel_app.Visible = False
                excel_app.DisplayAlerts = False

                workbook = excel_app.Workbooks.Open(
                    os.path.abspath(file_path), ReadOnly=True
                )
                sheets = [str(ws.Name) for ws in workbook.Worksheets]
                workbook.Close(SaveChanges=False)
                excel_app.Application.Quit()

                self._log_info(f"✅ COM Excel: {len(sheets)} sheets found")
                return sheets

            finally:
                pythoncom.CoUninitialize()

        except Exception as e:
            logger.debug(f"❌ COM Excel failed: {e}")

        return []

    def _read_excel_sheet(self, file_path: str, sheet_name: str) -> str:
        """Read a specific Excel sheet and convert to text

        Args:
            file_path: Path to Excel file
            sheet_name: Name of sheet to read

        Returns:
            str: Sheet content as text
        """
        try:
            import pandas as pd

            # Try different engines
            for engine in ["openpyxl", "xlrd", None]:
                try:
                    df = pd.read_excel(file_path, sheet_name=sheet_name, engine=engine)

                    # Convert to text format
                    content = ""

                    # Add column headers
                    if not df.empty:
                        headers = [str(col) for col in df.columns]
                        content += " | ".join(headers) + "\n"
                        content += "-" * len(" | ".join(headers)) + "\n"

                        # Add data rows
                        for _, row in df.iterrows():
                            row_data = [
                                str(val) if pd.notna(val) else "" for val in row.values
                            ]
                            content += " | ".join(row_data) + "\n"

                    return content.strip()

                except Exception:
                    continue

        except ImportError:
            pass

        raise Exception(f"Could not read sheet '{sheet_name}'")

    def _read_powerpoint(self, file_path: str) -> Dict[str, Any]:
        """Read PowerPoint files with multiple fallback methods

        Args:
            file_path: Path to PowerPoint file

        Returns:
            Dict[str, Any]: PowerPoint content and metadata
        """
        self._log_info(f"📊 Reading PowerPoint: {file_path}")

        file_extension = Path(file_path).suffix.lower()

        # Method 1: python-pptx (for .pptx files)
        if file_extension == ".pptx":
            try:
                from pptx import Presentation

                prs = Presentation(file_path)
                content = ""
                slide_count = 0

                for slide_num, slide in enumerate(prs.slides):
                    content += f"\n--- Slide {slide_num + 1} ---\n"

                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            content += shape.text + "\n"
                    slide_count += 1

                self._log_info(f"✅ python-pptx: {slide_count} slides processed")
                return {
                    "content": content.strip(),
                    "slide_count": slide_count,
                    "extraction_method": "python-pptx",
                }

            except ImportError:
                logger.warning("⚠️ python-pptx not available")
            except Exception as e:
                logger.warning(f"⚠️ python-pptx failed: {e}")

        # Method 2: Windows COM (for both .ppt and .pptx)
        try:
            import win32com.client
            import pythoncom

            pythoncom.CoInitialize()
            try:
                ppt_app = win32com.client.DispatchEx("PowerPoint.Application")
                ppt_app.Visible = False

                presentation = ppt_app.Presentations.Open(
                    os.path.abspath(file_path), ReadOnly=True
                )
                content = ""
                slide_count = presentation.Slides.Count

                for slide_num in range(1, slide_count + 1):
                    slide = presentation.Slides(slide_num)
                    content += f"\n--- Slide {slide_num} ---\n"

                    for shape in slide.Shapes:
                        if hasattr(shape, "TextFrame") and shape.TextFrame.HasText:
                            content += shape.TextFrame.TextRange.Text + "\n"

                presentation.Close()
                ppt_app.Application.Quit()

                self._log_info(f"✅ COM PowerPoint: {slide_count} slides processed")
                return {
                    "content": content.strip(),
                    "slide_count": slide_count,
                    "extraction_method": "COM PowerPoint",
                }

            finally:
                pythoncom.CoUninitialize()

        except ImportError:
            logger.warning("⚠️ pywin32 not available")
        except Exception as e:
            logger.warning(f"⚠️ COM PowerPoint failed: {e}")

        raise Exception("All PowerPoint reading methods failed")

    def _read_text(self, file_path: str) -> Dict[str, Any]:
        """Read text files with encoding detection

        Args:
            file_path: Path to text file

        Returns:
            Dict[str, Any]: Text content and metadata
        """
        logger.info(f"📝 Reading text file: {file_path}")

        # Detect encoding
        encoding = "utf-8"
        try:
            with open(file_path, "rb") as f:
                raw_data = f.read()
                detected = chardet.detect(raw_data)
                if detected["encoding"]:
                    encoding = detected["encoding"]
                    logger.info(f"🔍 Detected encoding: {encoding}")
        except:
            logger.warning("⚠️ Could not detect encoding, using utf-8")

        # Read content
        try:
            with open(file_path, "r", encoding=encoding) as f:
                content = f.read()

            lines_count = len(content.splitlines())
            word_count = len(content.split())

            self._log_info(f"✅ Text file: {lines_count} lines, {word_count} words")
            return {
                "content": content,
                "encoding": encoding,
                "lines_count": lines_count,
                "word_count": word_count,
                "extraction_method": "direct text reading",
            }

        except UnicodeDecodeError:
            # Fallback to different encodings
            for fallback_encoding in ["latin-1", "cp1252", "iso-8859-1"]:
                try:
                    with open(file_path, "r", encoding=fallback_encoding) as f:
                        content = f.read()

                    self._log_info(f"✅ Text file (fallback {fallback_encoding})")
                    return {
                        "content": content,
                        "encoding": fallback_encoding,
                        "extraction_method": f"text reading ({fallback_encoding})",
                    }
                except:
                    continue

            raise Exception("Could not read text file with any encoding")

    def get_supported_formats(self) -> Dict[str, List[str]]:
        """Get list of supported file formats

        Returns:
            Dict[str, List[str]]: Supported formats by category
        """
        return self.supported_formats.copy()

    def test_capabilities(self) -> Dict[str, Any]:
        """Test reading capabilities and available libraries

        Returns:
            Dict[str, Any]: Capability test results
        """
        logger.info("🧪 Testing file reading capabilities...")

        capabilities = {
            "pdf_support": False,
            "word_support": False,
            "excel_support": False,
            "powerpoint_support": False,
            "text_support": True,  # Always supported
            "available_libraries": [],
            "missing_libraries": [],
            "test_date": datetime.now().isoformat(),
        }

        # Test PDF support
        pdf_libs = ["PyPDF2", "fitz"]
        for lib in pdf_libs:
            try:
                __import__(lib)
                capabilities["available_libraries"].append(lib)
                capabilities["pdf_support"] = True
            except ImportError:
                capabilities["missing_libraries"].append(lib)

        # Test Word support
        word_libs = ["docx", "docx2txt", "win32com.client"]
        for lib in word_libs:
            try:
                __import__(lib)
                capabilities["available_libraries"].append(lib)
                capabilities["word_support"] = True
            except ImportError:
                capabilities["missing_libraries"].append(lib)

        # Test Excel support
        excel_libs = ["pandas", "openpyxl", "xlrd", "pyxlsb", "win32com.client"]
        for lib in excel_libs:
            try:
                __import__(lib)
                if lib not in capabilities["available_libraries"]:
                    capabilities["available_libraries"].append(lib)
                capabilities["excel_support"] = True
            except ImportError:
                if lib not in capabilities["missing_libraries"]:
                    capabilities["missing_libraries"].append(lib)

        # Test PowerPoint support
        ppt_libs = ["pptx", "win32com.client"]
        for lib in ppt_libs:
            try:
                __import__(lib)
                if lib not in capabilities["available_libraries"]:
                    capabilities["available_libraries"].append(lib)
                capabilities["powerpoint_support"] = True
            except ImportError:
                if lib not in capabilities["missing_libraries"]:
                    capabilities["missing_libraries"].append(lib)

        # Summary
        total_formats = len(self.supported_formats)
        supported_formats = sum(
            [
                capabilities["pdf_support"],
                capabilities["word_support"],
                capabilities["excel_support"],
                capabilities["powerpoint_support"],
                capabilities["text_support"],
            ]
        )

        logger.info(
            f"📊 Capability test complete: {supported_formats}/{total_formats} format types supported"
        )

        return capabilities


# Fire CLI Functions


def read_file(file_path: str) -> None:
    """Read and display content from any supported file format

    Args:
        file_path: Path to the file to read
    """
    try:
        reader = UniversalFileReader()
        result = reader.read_file(file_path)

        if result["status"] == "success":
            logger.info(f"\n📄 File: {result['file_name']}")
            logger.info(f"📊 Size: {result['file_size_mb']} MB")
            logger.info(f"🔧 Method: {result.get('extraction_method', 'unknown')}")
            logger.info(f"📝 Content length: {len(result['content'])} characters")
            logger.info("\n" + "=" * 60)

            # Display content (truncated if very long)
            content = result["content"]
            if len(content) > 2000:
                logger.info(content[:2000] + "\n\n... (content truncated)")
                logger.info(f"[Showing first 2000 of {len(content)} characters]")
            else:
                logger.info(content)

            logger.info("=" * 60)
        else:
            logger.info(f"❌ Error reading file: {result['error']}")

    except Exception as e:
        logger.info(f"❌ Failed to read file: {e}")


def get_file_info(file_path: str) -> None:
    """Get detailed information about a file

    Args:
        file_path: Path to the file
    """
    try:
        reader = UniversalFileReader()
        result = reader.read_file(file_path)

        logger.info(f"\n📋 File Information:")
        logger.info(f"Name: {result['file_name']}")
        logger.info(f"Path: {result['file_path']}")
        logger.info(f"Extension: {result['file_extension']}")
        logger.info(
            f"Size: {result['file_size_mb']} MB ({result['file_size_bytes']:,} bytes)"
        )
        logger.info(f"Modified: {result['modified_time']}")
        logger.info(f"Status: {result['status']}")

        if result["status"] == "success":
            logger.info(
                f"Extraction method: {result.get('extraction_method', 'unknown')}"
            )
            logger.info(f"Content length: {len(result['content'])} characters")

            # Show format-specific info
            if "total_pages" in result:
                logger.info(f"Total pages: {result['total_pages']}")
            if "total_sheets" in result:
                logger.info(f"Total sheets: {result['total_sheets']}")
                logger.info(f"Sheet names: {', '.join(result['sheet_names'])}")
            if "slide_count" in result:
                logger.info(f"Total slides: {result['slide_count']}")
        else:
            logger.info(f"Error: {result['error']}")

    except Exception as e:
        logger.info(f"❌ Failed to get file info: {e}")


def test_capabilities() -> None:
    """Test and display file reading capabilities"""
    try:
        reader = UniversalFileReader()
        capabilities = reader.test_capabilities()

        logger.info(f"🧪 File Reading Capabilities Test")
        logger.info("=" * 50)
        logger.info(f"📄 PDF Support: {'✅' if capabilities['pdf_support'] else '❌'}")
        logger.info(
            f"📝 Word Support: {'✅' if capabilities['word_support'] else '❌'}"
        )
        logger.info(
            f"📊 Excel Support: {'✅' if capabilities['excel_support'] else '❌'}"
        )
        logger.info(
            f"📺 PowerPoint Support: {'✅' if capabilities['powerpoint_support'] else '❌'}"
        )
        logger.info(
            f"📃 Text Support: {'✅' if capabilities['text_support'] else '❌'}"
        )

        logger.info(f"📚 Available Libraries:")
        for lib in capabilities["available_libraries"]:
            logger.info(f"  ✅ {lib}")

        if capabilities["missing_libraries"]:
            logger.info(f"⚠️ Missing Libraries:")
            for lib in capabilities["missing_libraries"]:
                logger.info(f"  ❌ {lib}")

        supported_formats = reader.get_supported_formats()
        logger.info(f"\n📋 Supported File Extensions:")
        for format_type, extensions in supported_formats.items():
            status = "✅" if capabilities.get(f"{format_type}_support", False) else "❌"
            logger.info(f"  {status} {format_type.title()}: {', '.join(extensions)}")

    except Exception as e:
        logger.info(f"❌ Capability test failed: {e}")


def list_supported_formats() -> None:
    """List all supported file formats"""
    reader = UniversalFileReader()
    supported_formats = reader.get_supported_formats()

    logger.info(f"\n📋 Supported File Formats:")
    logger.info("=" * 40)
    for format_type, extensions in supported_formats.items():
        logger.info(f"{format_type.title()}: {', '.join(extensions)}")

    logger.info(f"\nTotal: {len(reader.all_extensions)} supported extensions")


if __name__ == "__main__":
    fire.Fire()
